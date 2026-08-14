"""Deterministic text extraction and extraction-quality scoring."""

import json
import string
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.config import config
from app.modules.html_processing import html_to_text, sanitize_html
from app.schemas.ingestion_models import FileInspection, NativeExtraction


TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8-sig").strip()
    except UnicodeDecodeError:
        return path.read_text(encoding="cp1252").strip()


def _extract_docx(path: Path) -> str:
    document = Document(path)
    blocks: list[str] = []
    for paragraph in document.paragraphs:
        if text := paragraph.text.strip():
            blocks.append(text)
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("\n".join(rows))
    for section in document.sections:
        for container in (section.header, section.footer):
            for paragraph in container.paragraphs:
                if text := paragraph.text.strip():
                    blocks.append(text)
    return "\n\n".join(blocks).strip()


def _extract_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    blocks: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if value is None else str(value).strip() for value in row]
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
    finally:
        workbook.close()
    return "\n\n".join(blocks).strip()


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(path)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                if text := shape.text.strip():
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"Slide {index}\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()


def _printable_ratio(text: str) -> float:
    if not text:
        return 0.0
    readable = sum(character.isprintable() or character in string.whitespace for character in text)
    return readable / len(text)


def _score_text(text: str, *, text_native: bool) -> tuple[bool, float]:
    ratio = _printable_ratio(text)
    minimum = 1 if text_native else config.MIN_NATIVE_TEXT_CHARS
    acceptable = len(text.strip()) >= minimum and ratio >= config.MIN_PRINTABLE_RATIO
    return acceptable, ratio


def extract_native(inspection: FileInspection) -> NativeExtraction:
    path = Path(inspection.path)
    extension = inspection.extension
    if extension in IMAGE_EXTENSIONS:
        return NativeExtraction(text="", extractor="none", acceptable=False, printable_ratio=0.0)

    if extension == ".pdf":
        try:
            reader = PdfReader(path)
            if reader.is_encrypted:
                raise ValueError("Password-protected PDF files are not supported")
            page_texts = [(page.extract_text() or "").strip() for page in reader.pages]
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError(f"Cannot read PDF file: {exc}") from exc
        text = "\n\n".join(page for page in page_texts if page).strip()
        page_count = len(page_texts)
        useful_pages = sum(len(page) >= config.MIN_PAGE_TEXT_CHARS for page in page_texts)
        coverage = useful_pages / page_count if page_count else 0.0
        acceptable, ratio = _score_text(text, text_native=False)
        acceptable = acceptable and coverage >= config.MIN_TEXT_PAGE_COVERAGE
        partial = 0 < useful_pages < page_count
        warnings = []
        if partial:
            warnings.append(f"Only {useful_pages}/{page_count} PDF pages contain extractable text")
        return NativeExtraction(
            text=text,
            extractor="pypdf",
            printable_ratio=ratio,
            page_coverage=coverage,
            page_count=page_count,
            has_partial_pages=partial,
            acceptable=acceptable,
            warnings=warnings,
        )

    try:
        if extension in {".html", ".htm"}:
            raw = _read_text(path)
            text = html_to_text(sanitize_html(raw))
            extractor = "beautifulsoup"
        elif extension == ".json":
            raw = _read_text(path)
            text = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
            extractor = "json"
        elif extension in TEXT_EXTENSIONS:
            text = _read_text(path)
            extractor = "text"
        elif extension == ".docx":
            text = _extract_docx(path)
            extractor = "python-docx"
        elif extension == ".xlsx":
            text = _extract_xlsx(path)
            extractor = "openpyxl"
        elif extension == ".pptx":
            text = _extract_pptx(path)
            extractor = "python-pptx"
        else:
            return NativeExtraction(text="", extractor="none", acceptable=False)
    except Exception as exc:
        raise ValueError(f"Cannot extract {extension or 'file'} content: {exc}") from exc

    acceptable, ratio = _score_text(text, text_native=extension in TEXT_EXTENSIONS)
    return NativeExtraction(
        text=text,
        extractor=extractor,
        printable_ratio=ratio,
        acceptable=acceptable,
    )
