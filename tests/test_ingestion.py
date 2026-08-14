import asyncio
import json
from pathlib import Path

import httpx
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.config import config
from app.modules.file_inspection import inspect_file
from app.modules.gemini_converter import GeminiHtmlConverter
from app.modules.html_processing import html_to_text, prepare_html_chunks, sanitize_html
from app.modules.native_extraction import extract_native
from app.schemas.ingestion_models import FileInspection


def test_inspection_uses_signature_and_rejects_spoofed_extension(tmp_path):
    image = tmp_path / "document.pdf"
    image.write_bytes(b"\x89PNG\r\n\x1a\n" + b"payload")
    try:
        inspect_file(str(image), "document.pdf")
    except ValueError as exc:
        assert "image/png" in str(exc)
    else:
        raise AssertionError("Spoofed extension was accepted")


def test_sanitize_and_structure_aware_chunking():
    unsafe = """
    ```html
    <article onclick="steal()"><script>alert(1)</script>
      <section data-page="1"><h1>Overview</h1><p>Safe paragraph.</p></section>
      <section data-page="2"><h2>Numbers</h2><table><tr><th>A</th><th>B</th></tr>
      <tr><td>1</td><td>2</td></tr></table></section>
    </article>
    ```
    """
    canonical = sanitize_html(unsafe)
    assert "script" not in canonical
    assert "onclick" not in canonical
    assert "Safe paragraph" in html_to_text(canonical)
    chunks = prepare_html_chunks(
        canonical,
        "source.pdf",
        "source.pdf",
        document_id="abc123",
        metadata={"ingestion_route": "gemini_html"},
    )
    assert chunks
    assert chunks[0]["document_id"] == "abc123"
    assert chunks[0]["heading_path"] == ["Overview"]
    assert {item["page"] for item in chunks} == {"1", "2"}


def test_gemini_converter_calls_interactions_api_and_caches(tmp_path, monkeypatch):
    image = tmp_path / "scan.png"
    image.write_bytes(b"\x89PNG\r\n\x1a\n" + b"scan")
    inspection = FileInspection(
        path=str(image),
        original_name="scan.png",
        extension=".png",
        detected_mime="image/png",
        size_bytes=image.stat().st_size,
        sha256="a" * 64,
    )
    calls = 0

    def handler(request: httpx.Request) -> httpx.Response:
        nonlocal calls
        calls += 1
        assert request.headers["x-goog-api-key"] == "test-key"
        payload = json.loads(request.content)
        assert payload["model"] == "gemini-3.5-flash-lite"
        assert payload["input"][1]["type"] == "image"
        return httpx.Response(
            200,
            json={
                "steps": [
                    {"type": "model", "content": [{"type": "text", "text": "<article><h1>Scan</h1><p>Hello</p></article>"}]}
                ]
            },
        )

    monkeypatch.setattr(config, "NORMALIZED_DIR", str(tmp_path / "normalized"))
    converter = GeminiHtmlConverter(
        api_key="test-key",
        model="gemini-3.5-flash-lite",
        transport=httpx.MockTransport(handler),
    )
    first = asyncio.run(converter.convert(inspection))
    second = asyncio.run(converter.convert(inspection))
    assert first.cached is False
    assert second.cached is True
    assert calls == 1
    assert Path(first.artifact_path).exists()


def test_office_formats_use_native_extractors(tmp_path):
    repeated = "Native office document content " * 5

    docx_path = tmp_path / "sample.docx"
    docx = Document()
    docx.add_paragraph(repeated)
    table = docx.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Column A"
    table.cell(0, 1).text = "Column B"
    docx.save(docx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Description"])
    sheet.append(["Example", repeated])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    text_box = slide.shapes.add_textbox(0, 0, 5_000_000, 2_000_000)
    text_box.text = repeated
    presentation.save(pptx_path)

    expected = {
        docx_path: "python-docx",
        xlsx_path: "openpyxl",
        pptx_path: "python-pptx",
    }
    for path, extractor in expected.items():
        inspection = inspect_file(str(path), path.name)
        extraction = extract_native(inspection)
        assert extraction.acceptable
        assert extraction.extractor == extractor
        assert "Native office document content" in extraction.text
